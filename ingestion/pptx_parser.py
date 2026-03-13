"""
PPTX Parser — Layer 1: Ingestion
Extracts text, tables and speaker notes from PowerPoint pitch decks.
"""
import hashlib
from pathlib import Path

from pptx import Presentation
from pptx.util import Pt
from pptx.enum.shapes import MSO_SHAPE_TYPE

from ingestion.models import DeckDocument, Slide, SlideTable


def parse_pptx(file_path: str, scrub_pii: bool = False) -> DeckDocument:
    """
    Parse a PPTX pitch deck into a DeckDocument.

    Args:
        file_path: Absolute path to the PPTX file.
        scrub_pii: If True, run PII scrubber after parsing.

    Returns:
        DeckDocument with raw slides populated.
    """
    path = Path(file_path)
    if not path.exists():
        raise FileNotFoundError(f"PPTX not found: {file_path}")
    if path.suffix.lower() not in (".pptx", ".ppt"):
        raise ValueError(f"Expected .pptx file, got: {path.suffix}")

    file_hash = _hash_file(file_path)

    doc = DeckDocument(
        source_format="pptx",
        source_filename=path.name,
        metadata_hash=file_hash,
    )

    prs = Presentation(file_path)
    doc.slide_count = len(prs.slides)

    for slide_num, pptx_slide in enumerate(prs.slides, start=1):
        slide = Slide(slide_number=slide_num)

        text_parts = []
        image_count = 0

        for shape in pptx_slide.shapes:
            shape_type = shape.shape_type

            # ── Text frames ───────────────────────────────────────────────
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = " ".join(
                        run.text for run in para.runs if run.text.strip()
                    ).strip()
                    if line:
                        text_parts.append(line)

            # ── Tables ────────────────────────────────────────────────────
            elif shape.has_table:
                tbl = shape.table
                rows_data = []
                for row in tbl.rows:
                    row_data = []
                    for cell in row.cells:
                        cell_text = cell.text_frame.text.strip() if cell.text_frame else ""
                        row_data.append(cell_text)
                    rows_data.append(row_data)

                if rows_data:
                    # First row as header if it has content
                    if len(rows_data) >= 2 and any(rows_data[0]):
                        slide.tables.append(SlideTable(
                            headers=rows_data[0],
                            rows=rows_data[1:]
                        ))
                    else:
                        slide.tables.append(SlideTable(
                            headers=[],
                            rows=rows_data
                        ))

            # ── Images / charts ───────────────────────────────────────────
            elif shape_type == MSO_SHAPE_TYPE.PICTURE:
                image_count += 1
            elif hasattr(shape, "chart"):
                slide.has_charts = True

        slide.raw_text = "\n".join(text_parts).strip()
        slide.image_count = image_count

        # Heuristic chart detection: images with sparse text
        if image_count > 0 and len(slide.raw_text) < 80:
            slide.has_charts = True

        # ── Speaker notes ─────────────────────────────────────────────────
        if pptx_slide.has_notes_slide:
            notes_frame = pptx_slide.notes_slide.notes_text_frame
            if notes_frame:
                notes_text = notes_frame.text.strip()
                # Exclude PowerPoint's default placeholder text
                if notes_text and notes_text.lower() != "click to add notes":
                    slide.speaker_notes = notes_text

        doc.slides.append(slide)

    if scrub_pii:
        from cleaning.pii_scrubber import scrub_deck
        scrub_deck(doc)

    return doc


def _hash_file(file_path: str) -> str:
    sha256 = hashlib.sha256()
    with open(file_path, "rb") as f:
        for chunk in iter(lambda: f.read(8192), b""):
            sha256.update(chunk)
    return sha256.hexdigest()