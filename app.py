import os
import json
import sqlite3
import threading
import uuid
import shutil
from datetime import datetime
from pathlib import Path
from werkzeug.security import generate_password_hash, check_password_hash

from flask import Flask, request, jsonify, send_from_directory, Response, session, send_file
from flask_cors import CORS

app = Flask(__name__, static_folder="frontend", static_url_path="")
app.secret_key = os.environ.get("SECRET_KEY", "change_this_in_production_please")
CORS(app, supports_credentials=True)

BASE_DIR      = Path(__file__).parent
DB_PATH       = BASE_DIR / "dd_platform.db"
UPLOAD_DIR    = BASE_DIR / "uploads"
REPORT_DIR    = BASE_DIR / "reports"
SNAPSHOT_DIR  = BASE_DIR / "snapshots"
CHAT_FILE_DIR = BASE_DIR / "chat_files"

for d in [UPLOAD_DIR, REPORT_DIR, SNAPSHOT_DIR, CHAT_FILE_DIR]:
    d.mkdir(exist_ok=True)

ALLOWED_EXTENSIONS      = {".pdf", ".pptx", ".docx"}
CHAT_ALLOWED_EXTENSIONS = {".pdf", ".pptx", ".docx", ".txt", ".csv"}

# Admin password — change this before going live
ADMIN_PASSWORD = os.environ.get("ADMIN_PASSWORD", "startupscale2026")


# ── DB ─────────────────────────────────────────────────────────────────────────

def get_db():
    conn = sqlite3.connect(DB_PATH)
    conn.row_factory = sqlite3.Row
    return conn


def init_db():
    with get_db() as db:
        db.executescript("""
        CREATE TABLE IF NOT EXISTS users (
            id            INTEGER PRIMARY KEY AUTOINCREMENT,
            full_name     TEXT,
            email         TEXT UNIQUE,
            company       TEXT,
            role          TEXT,
            password_hash TEXT,
            created_at    TEXT DEFAULT (datetime('now'))
        );

        CREATE TABLE IF NOT EXISTS analyses (
            id            TEXT PRIMARY KEY,
            user_email    TEXT,
            filename      TEXT,
            startup_name  TEXT,
            slide_count   INTEGER DEFAULT 0,
            total_checks  INTEGER DEFAULT 0,
            flagged_count INTEGER DEFAULT 0,
            unclear_count INTEGER DEFAULT 0,
            ic_memo_risks INTEGER DEFAULT 0,
            deal_score    INTEGER DEFAULT 0,
            version       INTEGER DEFAULT 1,
            status        TEXT DEFAULT 'running',
            created_at    TEXT
        );

        CREATE TABLE IF NOT EXISTS findings (
            id          INTEGER PRIMARY KEY AUTOINCREMENT,
            analysis_id TEXT,
            version     INTEGER DEFAULT 1,
            check_id    TEXT,
            category    TEXT,
            severity    TEXT,
            status      TEXT,
            label       TEXT,
            evidence    TEXT DEFAULT ''
        );

        CREATE TABLE IF NOT EXISTS iterations (
            id          INTEGER PRIMARY KEY AUTOINCREMENT,
            analysis_id TEXT,
            version     INTEGER DEFAULT 1,
            role        TEXT,
            message     TEXT,
            attachment  TEXT DEFAULT '',
            created_at  TEXT
        );

        CREATE TABLE IF NOT EXISTS reports (
            id           INTEGER PRIMARY KEY AUTOINCREMENT,
            analysis_id  TEXT,
            version      INTEGER DEFAULT 1,
            report_type  TEXT,
            file_path    TEXT,
            created_at   TEXT
        );

        CREATE TABLE IF NOT EXISTS audit_log (
            id          INTEGER PRIMARY KEY AUTOINCREMENT,
            user_email  TEXT,
            action      TEXT,
            detail      TEXT DEFAULT '',
            ip_address  TEXT DEFAULT '',
            created_at  TEXT DEFAULT (datetime('now'))
        );

        CREATE TABLE IF NOT EXISTS snapshots (
            id          INTEGER PRIMARY KEY AUTOINCREMENT,
            analysis_id TEXT,
            version     INTEGER,
            snapshot    TEXT,
            created_at  TEXT DEFAULT (datetime('now'))
        );
        """)


# ── Audit log helper ───────────────────────────────────────────────────────────

def audit(user_email, action, detail=""):
    ip = request.remote_addr if request else "system"
    try:
        with get_db() as db:
            db.execute(
                "INSERT INTO audit_log (user_email,action,detail,ip_address) VALUES (?,?,?,?)",
                (user_email or "anonymous", action, detail[:500], ip)
            )
    except Exception:
        pass


# ── Snapshot helper ────────────────────────────────────────────────────────────

def save_snapshot(aid, version, findings):
    """Save current analysis state as a JSON snapshot for rollback."""
    try:
        snap = json.dumps({
            "analysis_id": aid,
            "version": version,
            "findings": findings,
            "saved_at": datetime.now().isoformat()
        })
        with get_db() as db:
            db.execute(
                "INSERT INTO snapshots (analysis_id,version,snapshot) VALUES (?,?,?)",
                (aid, version, snap)
            )
        # Also write to disk
        snap_path = SNAPSHOT_DIR / f"{aid}_v{version}.json"
        snap_path.write_text(snap, encoding="utf-8")
    except Exception as e:
        print(f"[SNAPSHOT] Failed to save snapshot: {e}")


# ── Deal score calculator ──────────────────────────────────────────────────────

def calculate_deal_score(findings):
    """
    Calculate a deal score out of 100 based on findings.
    Uses percentage-based scoring:
    - CLEAR checks count as full credit
    - UNCLEAR checks count as half credit  
    - FLAGGED checks count as zero
    Score = (CLEAR + UNCLEAR×0.5) / total_checks × 100
    Returns score and dimension scores for radar chart.
    """
    category_map = {
        "MS":    "Structure",
        "MKT":   "Market",
        "TECH":  "Product",
        "FIN":   "Financials",
        "GTM":   "Traction",
        "COMP":  "Competition",
        "TEAM":  "Team",
        "LEGAL": "Risk",
        "IC":    "Risk",
    }

    # Track per-dimension counts
    dim_clear   = {v: 0 for v in set(category_map.values())}
    dim_unclear = {v: 0 for v in set(category_map.values())}
    dim_total   = {v: 0 for v in set(category_map.values())}

    total_clear   = 0
    total_unclear = 0
    total_checks  = 0

    for f in findings:
        status   = f.get("status", "")
        category = f.get("category", "IC")
        dim      = category_map.get(category, "Risk")

        total_checks += 1
        dim_total[dim] = dim_total.get(dim, 0) + 1

        if status == "CLEAR":
            total_clear += 1
            dim_clear[dim] = dim_clear.get(dim, 0) + 1
        elif status == "UNCLEAR":
            total_unclear += 1
            dim_unclear[dim] = dim_unclear.get(dim, 0) + 1

    # Overall score
    if total_checks == 0:
        total_score = 0
    else:
        earned = total_clear + (total_unclear * 0.5)
        total_score = round((earned / total_checks) * 100)

    total_score = max(0, min(100, total_score))

    # Per-dimension scores (0-10)
    radar = {}
    for dim in set(category_map.values()):
        t = dim_total.get(dim, 0)
        if t == 0:
            radar[dim] = 10.0  # no checks in this dimension = full score
        else:
            c = dim_clear.get(dim, 0)
            u = dim_unclear.get(dim, 0)
            dim_score = (c + u * 0.5) / t * 10
            radar[dim] = round(dim_score, 1)

    return total_score, radar

def calculate_section_scores(findings):
    """
    Calculate 0-10 score per section category.
    Uses same deduction logic as deal score but per category.
    """
    category_map = {
        "MS":    "Structure",
        "MKT":   "Market",
        "TECH":  "Technology",
        "FIN":   "Financials",
        "GTM":   "Go-to-Market",
        "COMP":  "Competition",
        "TEAM":  "Team",
        "LEGAL": "Legal & Risk",
        "IC":    "IC Checks",
    }

    deductions = {"HIGH": 8, "MEDIUM": 4, "LOW": 1}
    unclear_deductions = {"HIGH": 4, "MEDIUM": 2, "LOW": 1}

    # Track scores and check counts per category
    scores = {v: 100 for v in category_map.values()}
    counts = {v: 0 for v in category_map.values()}

    for f in findings:
        status   = f.get("status", "")
        severity = f.get("severity", "LOW")
        category = f.get("category", "IC")
        dim      = category_map.get(category, "IC Checks")

        if dim not in scores:
            continue

        counts[dim] += 1

        if status == "FLAGGED":
            ded = deductions.get(severity, 1)
            scores[dim] = max(0, scores[dim] - ded * 2)
        elif status == "UNCLEAR":
            ded = unclear_deductions.get(severity, 1)
            scores[dim] = max(0, scores[dim] - ded)

    # Only include categories that had checks run
    result = {}
    for cat, score in scores.items():
        if counts[cat] > 0:
            result[cat] = {
                "score": round(score / 10, 1),
                "checks": counts[cat],
                "color": "green" if score >= 70 else "amber" if score >= 40 else "red"
            }

    return result

def score_verdict(score):
    if score >= 75:
        return "Proceed to Full Diligence"
    elif score >= 55:
        return "Conditional Proceed"
    elif score >= 35:
        return "Pass — Conditions Apply"
    else:
        return "Not Investment Ready"


def score_color(score):
    if score >= 75:
        return "green"
    elif score >= 55:
        return "amber"
    elif score >= 35:
        return "amber"
    else:
        return "red"

# ── Progress store ─────────────────────────────────────────────────────────────

_progress = {}


def push(aid, step, label, pct, done=False, discovery=None):
    event = {"step": step, "label": label, "pct": pct, "done": done}
    if discovery:
        event["discovery"] = discovery
    _progress.setdefault(aid, []).append(json.dumps(event))


# ── AUTH ───────────────────────────────────────────────────────────────────────

@app.route("/api/register", methods=["POST"])
def register():
    data      = request.json or {}
    full_name = (data.get("full_name") or "").strip()
    email     = (data.get("email") or "").strip().lower()
    company   = (data.get("company") or "").strip()
    role      = (data.get("role") or "investor").strip()
    password  = (data.get("password") or "").strip()

    if not all([full_name, email, company, password]):
        return jsonify({"error": "All fields required"}), 400

    with get_db() as db:
        if db.execute("SELECT id FROM users WHERE email=?", (email,)).fetchone():
            return jsonify({"error": "Email already registered"}), 400
        db.execute(
            "INSERT INTO users (full_name,email,company,role,password_hash) VALUES (?,?,?,?,?)",
            (full_name, email, company, role, generate_password_hash(password))
        )

    session["user"] = email
    audit(email, "REGISTER", f"New {role} account — {company}")
    return jsonify({"ok": True})


@app.route("/api/login", methods=["POST"])
def login():
    data     = request.json or {}
    email    = (data.get("email") or "").strip().lower()
    password = (data.get("password") or "").strip()

    with get_db() as db:
        user = db.execute("SELECT * FROM users WHERE email=?", (email,)).fetchone()

    if not user or not check_password_hash(user["password_hash"], password):
        audit(email, "LOGIN_FAILED", "Invalid credentials")
        return jsonify({"error": "Invalid email or password"}), 401

    session["user"] = email
    audit(email, "LOGIN", f"Logged in from {request.remote_addr}")
    return jsonify({"ok": True})


@app.route("/api/me")
def me():
    if "user" not in session:
        return jsonify({"error": "Not logged in"}), 401
    with get_db() as db:
        user = db.execute("SELECT * FROM users WHERE email=?", (session["user"],)).fetchone()
    if not user:
        return jsonify({"error": "User not found"}), 404
    u = dict(user)
    u.pop("password_hash", None)
    return jsonify(u)


@app.route("/api/logout", methods=["POST"])
def logout():
    audit(session.get("user"), "LOGOUT")
    session.clear()
    return jsonify({"ok": True})


# ── ADMIN AUTH ─────────────────────────────────────────────────────────────────

@app.route("/api/admin/login", methods=["POST"])
def admin_login():
    data = request.json or {}
    if data.get("password") == ADMIN_PASSWORD:
        session["admin"] = True
        audit("admin", "ADMIN_LOGIN", f"Admin login from {request.remote_addr}")
        return jsonify({"ok": True})
    return jsonify({"error": "Invalid admin password"}), 401


def admin_required(fn):
    from functools import wraps
    @wraps(fn)
    def wrapper(*args, **kwargs):
        if not session.get("admin"):
            return jsonify({"error": "Admin access required"}), 403
        return fn(*args, **kwargs)
    return wrapper


# ── ADMIN ROUTES ───────────────────────────────────────────────────────────────

@app.route("/api/admin/stats")
@admin_required
def admin_stats():
    with get_db() as db:
        total_users     = db.execute("SELECT COUNT(*) FROM users").fetchone()[0]
        total_analyses  = db.execute("SELECT COUNT(*) FROM analyses").fetchone()[0]
        complete        = db.execute("SELECT COUNT(*) FROM analyses WHERE status='complete'").fetchone()[0]
        failed          = db.execute("SELECT COUNT(*) FROM analyses WHERE status='failed'").fetchone()[0]
        total_reports   = db.execute("SELECT COUNT(*) FROM reports").fetchone()[0]
        avg_score       = db.execute("SELECT AVG(deal_score) FROM analyses WHERE status='complete'").fetchone()[0]
        avg_flags       = db.execute("SELECT AVG(flagged_count) FROM analyses WHERE status='complete'").fetchone()[0]

        # Analyses per day (last 14 days)
        daily = db.execute("""
            SELECT DATE(created_at) as day, COUNT(*) as count
            FROM analyses
            WHERE created_at >= DATE('now', '-14 days')
            GROUP BY DATE(created_at)
            ORDER BY day
        """).fetchall()

        # Most common flags
        top_flags = db.execute("""
            SELECT check_id, COUNT(*) as freq
            FROM findings WHERE status='FLAGGED'
            GROUP BY check_id ORDER BY freq DESC LIMIT 8
        """).fetchall()

    return jsonify({
        "total_users":    total_users,
        "total_analyses": total_analyses,
        "complete":       complete,
        "failed":         failed,
        "total_reports":  total_reports,
        "avg_score":      round(avg_score or 0, 1),
        "avg_flags":      round(avg_flags or 0, 1),
        "daily_analyses": [dict(r) for r in daily],
        "top_flags":      [dict(r) for r in top_flags],
    })


@app.route("/api/admin/users")
@admin_required
def admin_users():
    with get_db() as db:
        rows = db.execute(
            "SELECT id,full_name,email,company,role,created_at FROM users ORDER BY created_at DESC"
        ).fetchall()
    return jsonify([dict(r) for r in rows])


@app.route("/api/admin/analyses")
@admin_required
def admin_analyses():
    with get_db() as db:
        rows = db.execute("""
            SELECT a.id, a.user_email, a.startup_name, a.filename,
                   a.slide_count, a.flagged_count, a.unclear_count,
                   a.deal_score, a.status, a.version, a.created_at,
                   u.full_name, u.company
            FROM analyses a
            LEFT JOIN users u ON a.user_email = u.email
            ORDER BY a.created_at DESC LIMIT 100
        """).fetchall()
    return jsonify([dict(r) for r in rows])


@app.route("/api/admin/audit")
@admin_required
def admin_audit():
    with get_db() as db:
        rows = db.execute(
            "SELECT * FROM audit_log ORDER BY created_at DESC LIMIT 200"
        ).fetchall()
    return jsonify([dict(r) for r in rows])


@app.route("/api/admin/export")
@admin_required
def admin_export():
    """Export all data as Excel."""
    try:
        import openpyxl
        wb = openpyxl.Workbook()

        # Users sheet
        ws_users = wb.active
        ws_users.title = "Users"
        ws_users.append(["ID", "Name", "Email", "Company", "Role", "Joined"])
        with get_db() as db:
            for r in db.execute("SELECT id,full_name,email,company,role,created_at FROM users").fetchall():
                ws_users.append(list(r))

        # Analyses sheet
        ws_analyses = wb.create_sheet("Analyses")
        ws_analyses.append(["ID", "User", "Startup", "Slides", "Flagged", "Score", "Status", "Date"])
        with get_db() as db:
            for r in db.execute("SELECT id,user_email,startup_name,slide_count,flagged_count,deal_score,status,created_at FROM analyses").fetchall():
                ws_analyses.append(list(r))

        # Audit sheet
        ws_audit = wb.create_sheet("Audit Log")
        ws_audit.append(["ID", "User", "Action", "Detail", "IP", "Time"])
        with get_db() as db:
            for r in db.execute("SELECT * FROM audit_log ORDER BY created_at DESC").fetchall():
                ws_audit.append(list(r))

        export_path = BASE_DIR / "admin_export.xlsx"
        wb.save(str(export_path))
        return send_file(str(export_path), as_attachment=True,
                         download_name="startupscale360_export.xlsx")
    except ImportError:
        return jsonify({"error": "Install openpyxl: pip install openpyxl"}), 500


# ── RUN PIPELINE ───────────────────────────────────────────────────────────────

@app.route("/api/run", methods=["POST"])
def run():
    if "user" not in session:
        return jsonify({"error": "Unauthorized"}), 401

    if "file" not in request.files:
        return jsonify({"error": "No file uploaded"}), 400

    file = request.files["file"]
    if not file or not file.filename:
        return jsonify({"error": "Invalid file"}), 400

    ext = Path(file.filename).suffix.lower()
    if ext not in ALLOWED_EXTENSIONS:
        return jsonify({"error": "File must be PDF, PPTX or DOCX"}), 400

    startup_name = (request.form.get("startup_name") or "").strip()
    user_email   = session["user"]
    provider = request.form.get("provider") or os.environ.get("LLM_PROVIDER", "openai")
    use_llm      = request.form.get("use_llm", "true").lower() == "true"

    aid      = str(uuid.uuid4())[:8]
    tmp_path = UPLOAD_DIR / f"{aid}{ext}"
    file.save(tmp_path)
    _progress[aid] = []

    with get_db() as db:
        db.execute(
            """INSERT INTO analyses
               (id,user_email,filename,startup_name,slide_count,total_checks,
                flagged_count,unclear_count,ic_memo_risks,deal_score,version,status,created_at)
               VALUES (?,?,?,?,0,0,0,0,0,0,1,'running',?)""",
            (aid, user_email, file.filename,
             startup_name or file.filename, datetime.now().isoformat())
        )

    audit(user_email, "ANALYSIS_STARTED", f"{startup_name or file.filename}")

    def process(email_snap):
        try:
            push(aid, 1, "Ingesting deck...", 10)

            from pipeline import run_pipeline
            push(aid, 2, "Cleaning & classifying sections...", 28)

            result = run_pipeline(
                str(tmp_path),
                run_llm_pass=use_llm,
                llm_provider=provider
            )

            # Discovery event — sections found
            sections = sorted(result.deck.sections_present())
            push(aid, 3, "Running anomaly checks...", 55,
                 discovery=f"Found {len(sections)} sections: {', '.join(sections)}")

            push(aid, 4, "Generating DD report & IC memo...", 78)

            import threading as _threading
            _report_done = _threading.Event()
            _report_paths = {}
            _report_error = [None]

            _SECTION_LABELS = {
                "product_description": "Product Description",
                "technology_solution": "Technology Solution",
                "business_model": "Business Model",
                "technology_overview": "Technology Architecture",
                "comments_technology": "Technology Q&A",
                "team": "Team",
                "cap_table": "Cap Table & Funding",
                "traction": "Traction",
                "market_opportunity": "Market Opportunity",
                "competitive_landscape": "Competitive Landscape",
                "investment_rationale": "Investment Rationale",
                "areas_to_watch": "Areas to Watch",
                "final_call_a": "Final Recommendation",
                "executive_summary_b": "IC Executive Summary",
                "ic_summary_positive": "IC Positive Signals",
                "ic_summary_risks": "IC Key Risks",
                "ic_recommendation": "IC Recommendation",
            }

            def _on_section(section_key):
                label = _SECTION_LABELS.get(section_key, section_key.replace("_", " ").title())
                push(aid, 4, "Generating DD report & IC memo...", _heartbeat_pct,
                     discovery=f"✅ Written: {label}")

            def _generate():
                try:
                    from report_1.report_generator import generate_reports
                    _report_paths.update(generate_reports(
                        result,
                        output_dir=str(REPORT_DIR),
                        on_section_complete=_on_section
                    ))
                except Exception as e:
                    _report_error[0] = e
                finally:
                    _report_done.set()

            _threading.Thread(target=_generate, daemon=True).start()

            _heartbeat_pct = 78
            while not _report_done.wait(timeout=15):
                _heartbeat_pct = min(_heartbeat_pct + 2, 97)
                push(aid, 4, "Generating DD report & IC memo...", _heartbeat_pct)

            if _report_error[0]:
                raise _report_error[0]

            paths = _report_paths

            # ── Company Mode: generate founder-facing suggestions ──────────
            with get_db() as db:
                user_row = db.execute("SELECT role FROM users WHERE email=?", (email_snap,)).fetchone()
            user_role = (user_row["role"] if user_row else "investor") or "investor"

            suggestions_by_check_id = {}
            if user_role == "startup":
                push(aid, 5, "Generating improvement suggestions...", 92)
                from report_1.section_writer import SectionWriter
                checks_by_id = {c.id: c for c in result.config.anomalies}
                suggestion_writer = SectionWriter(provider=provider)

                from concurrent.futures import ThreadPoolExecutor as _TPE, as_completed as _ac

                flagged_unclear = [f for f in result.all_findings if f.status in ("FLAGGED", "UNCLEAR")]

                def _make_suggestion(f):
                    check = checks_by_id.get(f.anomaly_id)
                    slide_text = result.deck.text_for_sections(check.target_sections) if check else ""
                    return f.anomaly_id, suggestion_writer.write_suggestion(
                        f, [slide_text] if slide_text else [], []
                    )

                with _TPE(max_workers=5) as ex:
                    futures = {ex.submit(_make_suggestion, f): f for f in flagged_unclear}
                    for future in _ac(futures):
                        try:
                            check_id, suggestion = future.result()
                            suggestions_by_check_id[check_id] = suggestion
                        except Exception as e:
                            print(f"[SUGGESTION] Failed: {e}")

            push(aid, 5, "Calculating deal score...", 90)

            flagged  = len(result.flagged())
            unclear  = len(result.unclear())
            ic_risks = len(result.ic_memo_risks())
            total    = len(result.all_findings)
            version  = 1

            # Build findings list for score calculation
            findings_list = [
                {"status": f.status, "severity": f.severity, "category": f.category}
                for f in result.all_findings
            ]
            deal_score, radar = calculate_deal_score(findings_list)

            with get_db() as db:
                db.execute(
                    """UPDATE analyses SET
                       slide_count=?,total_checks=?,flagged_count=?,unclear_count=?,
                       ic_memo_risks=?,deal_score=?,status='complete' WHERE id=?""",
                    (result.deck.slide_count, total, flagged, unclear,
                     ic_risks, deal_score, aid)
                )

                findings_for_snapshot = []
                for f in result.all_findings:
                    suggestion_text = suggestions_by_check_id.get(f.anomaly_id, "")
                    db.execute(
                        """INSERT INTO findings
                           (analysis_id,version,check_id,category,severity,status,label,evidence,suggestion)
                           VALUES (?,?,?,?,?,?,?,?,?)""",
                        (aid, version, f.anomaly_id, f.category, f.severity,
                         f.status, f.label, f.evidence or "", suggestion_text)
                    )
                    findings_for_snapshot.append({
                        "check_id": f.anomaly_id, "category": f.category,
                        "severity": f.severity, "status": f.status,
                        "label": f.label, "evidence": f.evidence or "",
                        "suggestion": suggestion_text
                    })

                db.execute(
                    "INSERT INTO reports (analysis_id,version,report_type,file_path,created_at) VALUES (?,?,?,?,?)",
                    (aid, version, "framework_a", paths["framework_a_path"], datetime.now().isoformat())
                )
                db.execute(
                    "INSERT INTO reports (analysis_id,version,report_type,file_path,created_at) VALUES (?,?,?,?,?)",
                    (aid, version, "framework_b", paths["framework_b_path"], datetime.now().isoformat())
                )

                db.execute(
                    "INSERT INTO iterations (analysis_id,version,role,message,created_at) VALUES (?,?,?,?,?)",
                    (aid, version, "system",
                     f"Analysis complete. Deal score: {deal_score}/100. "
                     f"{flagged} item(s) flagged across {total} checks. "
                     f"Ask me anything about this deck or share new information to update the report.",
                     datetime.now().isoformat())
                )

            # Save snapshot
            save_snapshot(aid, version, findings_for_snapshot)

            audit(email_snap, "ANALYSIS_COMPLETE",
                  f"{aid} — score:{deal_score} flagged:{flagged}")
            push(aid, 5, "Complete", 100, True)

        except Exception as e:
                import traceback
                traceback.print_exc()
                with get_db() as db:
                    db.execute("UPDATE analyses SET status='failed' WHERE id=?", (aid,))
                audit(email_snap, "ANALYSIS_FAILED", f"{aid} — {str(e)[:200]}")
                push(aid, 5, f"Error: {str(e)[:120]}", 100, True)
        finally:
            if tmp_path.exists():
                tmp_path.unlink()

    threading.Thread(target=process, args=(user_email,), daemon=True).start()
    return jsonify({"analysis_id": aid})


# ── SSE PROGRESS ───────────────────────────────────────────────────────────────

@app.route("/api/progress/<aid>")
def progress(aid):
    def stream():
        import time
        sent = 0
        for _ in range(600):
            events = _progress.get(aid, [])
            while sent < len(events):
                yield f"data: {events[sent]}\n\n"
                sent += 1
            if events and json.loads(events[-1]).get("done"):
                break
            time.sleep(0.5)

    return Response(stream(), mimetype="text/event-stream",
                    headers={"Cache-Control": "no-cache", "X-Accel-Buffering": "no"})


# ── ANALYSIS RESULTS ───────────────────────────────────────────────────────────

@app.route("/api/analysis/<aid>")
def analysis(aid):
    if "user" not in session:
        return jsonify({"error": "Unauthorized"}), 401

    version = request.args.get("version", type=int)

    with get_db() as db:
        a = db.execute("SELECT * FROM analyses WHERE id=?", (aid,)).fetchone()
        if not a:
            return jsonify({"error": "Not found"}), 404

        # Data isolation — users only see their own analyses
        if a["user_email"] != session["user"]:
            return jsonify({"error": "Access denied"}), 403

        v = version or a["version"]

        findings = db.execute(
            "SELECT * FROM findings WHERE analysis_id=? AND version<=? ORDER BY id",
            (aid, v)
        ).fetchall()

        iterations = db.execute(
            "SELECT * FROM iterations WHERE analysis_id=? ORDER BY id",
            (aid,)
        ).fetchall()

        reports = db.execute(
            "SELECT * FROM reports WHERE analysis_id=? ORDER BY id",
            (aid,)
        ).fetchall()

        user_row = db.execute("SELECT role FROM users WHERE email=?", (a["user_email"],)).fetchone()
        user_role = (user_row["role"] if user_row else "investor") or "investor"

    findings_list = [dict(f) for f in findings]
    deal_score    = a["deal_score"] or 0
    _, radar      = calculate_deal_score(findings_list)

    # Build narrative
    flagged_high  = [f for f in findings_list if f["status"] == "FLAGGED" and f["severity"] == "HIGH"]
    flagged_med   = [f for f in findings_list if f["status"] == "FLAGGED" and f["severity"] == "MEDIUM"]
    clear_items   = [f for f in findings_list if f["status"] == "CLEAR"]
    top_risks     = (flagged_high + flagged_med)[:3]
    top_strengths = clear_items[:3]
    biggest_risk  = flagged_high[0] if flagged_high else (flagged_med[0] if flagged_med else None)

    # Overall benchmarking
    with get_db() as db:
        total_complete = db.execute(
            "SELECT COUNT(*) FROM analyses WHERE status='complete'"
        ).fetchone()[0]
        lower_score = db.execute(
            "SELECT COUNT(*) FROM analyses WHERE status='complete' AND deal_score<?",
            (deal_score,)
        ).fetchone()[0]

    percentile = round((lower_score / total_complete * 100)) if total_complete > 0 else 50

    # Section scores
    section_scores = calculate_section_scores(findings_list)

    # Platform average per section
    section_benchmarks = {}
    with get_db() as db:
        for cat_code, cat_name in [
            ("MS","Structure"), ("MKT","Market"), ("TECH","Technology"),
            ("FIN","Financials"), ("GTM","Go-to-Market"), ("COMP","Competition"),
            ("TEAM","Team"), ("LEGAL","Legal & Risk"), ("IC","IC Checks")
        ]:
            avg = db.execute("""
                SELECT AVG(CASE
                    WHEN f.status='FLAGGED' AND f.severity='HIGH'   THEN 0
                    WHEN f.status='FLAGGED' AND f.severity='MEDIUM' THEN 40
                    WHEN f.status='FLAGGED' AND f.severity='LOW'    THEN 70
                    WHEN f.status='UNCLEAR'                         THEN 60
                    ELSE 100
                END)
                FROM findings f
                JOIN analyses a ON f.analysis_id = a.id
                WHERE f.category=? AND a.status='complete'
            """, (cat_code,)).fetchone()[0]

            if avg is not None:
                section_benchmarks[cat_name] = round(avg / 10, 1)

    return jsonify({
        **dict(a),
        "user_role":          user_role,
        "findings":           findings_list,
        "iterations":         [dict(i) for i in iterations],
        "reports":            [dict(r) for r in reports],
        "deal_score":         deal_score,
        "score_verdict":      score_verdict(deal_score),
        "score_color":        score_color(deal_score),
        "radar":              radar,
        "top_risks":          top_risks,
        "top_strengths":      top_strengths,
        "biggest_risk":       dict(biggest_risk) if biggest_risk else None,
        "percentile":         percentile,
        "section_scores":     section_scores,
        "section_benchmarks": section_benchmarks,
        "sections": sorted(list(set(
            {"MS":"Structure","MKT":"Market","TECH":"Technology",
             "FIN":"Financials","GTM":"Go-to-Market","COMP":"Competition",
             "TEAM":"Team","LEGAL":"Legal","IC":"IC Checks"
            }.get(f["category"], f["category"])
            for f in findings_list))),
    })

# ── REPORT DOWNLOAD ────────────────────────────────────────────────────────────

@app.route("/api/reports/<int:report_id>/download")
def download_report(report_id):
    if "user" not in session:
        return jsonify({"error": "Unauthorized"}), 401

    with get_db() as db:
        rep = db.execute("SELECT * FROM reports WHERE id=?", (report_id,)).fetchone()
        if not rep:
            return jsonify({"error": "Report not found"}), 404
        # Check ownership
        a = db.execute("SELECT user_email FROM analyses WHERE id=?", (rep["analysis_id"],)).fetchone()
        if not a or a["user_email"] != session["user"]:
            return jsonify({"error": "Access denied"}), 403

    file_path = Path(rep["file_path"])
    if not file_path.exists():
        return jsonify({"error": "Report file missing from disk"}), 404

    audit(session["user"], "REPORT_DOWNLOAD", f"report_id:{report_id} type:{rep['report_type']}")
    return send_file(
        str(file_path), as_attachment=True, download_name=file_path.name,
        mimetype="application/vnd.openxmlformats-officedocument.wordprocessingml.document"
    )


# ── REPORT PREVIEW ─────────────────────────────────────────────────────────────

@app.route("/api/reports/<int:report_id>/preview")
def preview_report(report_id):
    if "user" not in session:
        return jsonify({"error": "Unauthorized"}), 401

    with get_db() as db:
        rep = db.execute("SELECT * FROM reports WHERE id=?", (report_id,)).fetchone()
        if not rep:
            return jsonify({"error": "Report not found"}), 404
        a = db.execute("SELECT user_email FROM analyses WHERE id=?", (rep["analysis_id"],)).fetchone()
        if not a or a["user_email"] != session["user"]:
            return jsonify({"error": "Access denied"}), 403

    file_path = Path(rep["file_path"])
    if not file_path.exists():
        return jsonify({"error": "Report file missing"}), 404

    try:
        from docx import Document
        doc = Document(str(file_path))
        # Extract with headings preserved
        sections = []
        for p in doc.paragraphs:
            if p.text.strip():
                if p.style.name.startswith("Heading"):
                    sections.append({"type": "heading", "text": p.text.strip()})
                else:
                    sections.append({"type": "paragraph", "text": p.text.strip()})
        preview = "\n\n".join(
            ("## " if s["type"] == "heading" else "") + s["text"]
            for s in sections[:40]
        )
    except Exception as e:
        preview = f"Preview unavailable: {str(e)}"

    return jsonify({"preview_text": preview, "report_type": rep["report_type"]})


# ── CHAT / ITERATE ─────────────────────────────────────────────────────────────

@app.route("/api/analysis/<aid>/iterate", methods=["POST"])
def iterate(aid):
    if "user" not in session:
        return jsonify({"error": "Unauthorized"}), 401

    # Handle both JSON and multipart (file upload)
    if request.content_type and "multipart" in request.content_type:
        message = (request.form.get("message") or "").strip()
        attachment_file = request.files.get("file")
    else:
        data = request.json or {}
        message = (data.get("message") or "").strip()
        attachment_file = None

    if not message and not attachment_file:
        return jsonify({"error": "Message or file required"}), 400

    attachment_name = ""
    attachment_text = ""

    # Handle attached file
    if attachment_file and attachment_file.filename:
        ext = Path(attachment_file.filename).suffix.lower()
        if ext in CHAT_ALLOWED_EXTENSIONS:
            fname = f"{aid}_{uuid.uuid4().hex[:6]}{ext}"
            fpath = CHAT_FILE_DIR / fname
            attachment_file.save(fpath)
            attachment_name = attachment_file.filename

            # Extract text from attached file
            try:
                if ext == ".pdf":
                    import pdfplumber
                    with pdfplumber.open(str(fpath)) as pdf:
                        attachment_text = "\n".join(
                            p.extract_text() or "" for p in pdf.pages[:5]
                        )
                elif ext in [".pptx", ".ppt"]:
                    from pptx import Presentation
                    prs = Presentation(str(fpath))
                    texts = []
                    for slide in list(prs.slides)[:5]:
                        for shape in slide.shapes:
                            if hasattr(shape, "text"):
                                texts.append(shape.text)
                    attachment_text = "\n".join(texts)
                elif ext == ".docx":
                    from docx import Document
                    doc = Document(str(fpath))
                    attachment_text = "\n".join(p.text for p in doc.paragraphs[:30])
                elif ext == ".txt":
                    attachment_text = fpath.read_text(encoding="utf-8")[:3000]
            except Exception as e:
                attachment_text = f"[Could not extract text: {e}]"

    with get_db() as db:
        a = db.execute("SELECT * FROM analyses WHERE id=?", (aid,)).fetchone()
        if not a:
            return jsonify({"error": "Analysis not found"}), 404
        if a["user_email"] != session["user"]:
            return jsonify({"error": "Access denied"}), 403

        findings = db.execute(
            "SELECT check_id,severity,status,label,evidence FROM findings WHERE analysis_id=?",
            (aid,)
        ).fetchall()

    flagged_summary = "\n".join([
            f"- [{f['severity']}] {f['label']}: {f['evidence'][:80]}"
            for f in findings if f["status"] == "FLAGGED"
        ]) or "No flagged items."

    unclear_summary = "\n".join([
        f"- {f['label']}: {f['evidence'][:60]}"
        for f in findings if f["status"] == "UNCLEAR"
    ]) or "None."

    clear_summary = "\n".join([
        f"- {f['label']}"
        for f in findings if f["status"] == "CLEAR"
    ]) or "None."

    deck_name = a["startup_name"] or a["filename"]
    current_v = a["version"]

    update_keywords = [
        "actually", "update", "correct", "new info", "they have", "confirmed",
        "just heard", "please note", "add", "the company", "i can confirm",
        "we have", "turns out", "clarification"
    ]
    is_update = any(kw in message.lower() for kw in update_keywords) or bool(attachment_name)

    # Build full context
    full_message = message
    if attachment_text:
        full_message += f"\n\n[Attached document: {attachment_name}]\n{attachment_text[:2000]}"

    system_prompt = f"""You are an AI due diligence assistant built into the StartupScale360 platform.
You have already run a full automated analysis on the pitch deck for "{deck_name}".

AUTOMATED PIPELINE FINDINGS — treat these as ground truth:

FLAGGED (red flags requiring attention):
{flagged_summary}

UNCLEAR (needs clarification):
{unclear_summary}

CLEAR (confirmed positives):
{clear_summary}

PLATFORM CAPABILITIES:
- The platform has already generated a DD Report and IC Memo — downloadable from the Download Reports section
- Users can upload additional documents to update the analysis
- The platform ran 44 checks across the full deck content

YOUR ROLE:
- You have full knowledge of all 44 check results above — use them to answer questions
- When asked to "deep dive" or "analyse further", use the evidence from the findings above — do NOT ask for the deck again
- When users upload new documents, explain which specific flags it resolves and why
- NEVER contradict the pipeline findings — if new info resolves a flag, acknowledge it but note verification is needed
- NEVER say you cannot generate documents — they are already generated
- Keep responses under 150 words, short paragraphs, no markdown headers
- Do not repeat the full findings list back — the user can already see them"""

    reply = ""
    try:
        from llm.adapter import get_llm
        llm = get_llm(os.environ.get("LLM_PROVIDER", "openai"))
        res   = llm.complete(system_prompt=system_prompt, user_prompt=full_message,
                             temperature=0.3, max_tokens=400)
        reply = res.text.strip()
    except Exception as e:
        reply = f"Noted. (AI response unavailable: {str(e)[:60]})"

    new_version = current_v + 1 if is_update else current_v

    with get_db() as db:
        db.execute(
            "INSERT INTO iterations (analysis_id,version,role,message,attachment,created_at) VALUES (?,?,?,?,?,?)",
            (aid, current_v, "user",
             message or f"Uploaded: {attachment_name}",
             attachment_name, datetime.now().isoformat())
        )
        db.execute(
            "INSERT INTO iterations (analysis_id,version,role,message,attachment,created_at) VALUES (?,?,?,?,?,?)",
            (aid, new_version, "assistant", reply, "", datetime.now().isoformat())
        )
        if is_update:
            db.execute("UPDATE analyses SET version=? WHERE id=?", (new_version, aid))

    audit(session["user"], "CHAT_MESSAGE",
          f"{aid} — {'update' if is_update else 'question'} v{new_version}")

    return jsonify({
        "mode":       "update" if is_update else "question",
        "reply":      reply,
        "version":    new_version,
        "attachment": attachment_name
    })


# ── SNAPSHOTS / ROLLBACK ───────────────────────────────────────────────────────

@app.route("/api/analysis/<aid>/snapshots")
def get_snapshots(aid):
    if "user" not in session:
        return jsonify({"error": "Unauthorized"}), 401
    with get_db() as db:
        a = db.execute("SELECT user_email FROM analyses WHERE id=?", (aid,)).fetchone()
        if not a or a["user_email"] != session["user"]:
            return jsonify({"error": "Access denied"}), 403
        snaps = db.execute(
            "SELECT id,version,created_at FROM snapshots WHERE analysis_id=? ORDER BY id DESC",
            (aid,)
        ).fetchall()
    return jsonify([dict(s) for s in snaps])


# ── HISTORY ────────────────────────────────────────────────────────────────────

@app.route("/api/history")
def history():
    if "user" not in session:
        return jsonify([])
    with get_db() as db:
        rows = db.execute(
            """SELECT * FROM analyses WHERE user_email=?
               ORDER BY created_at DESC LIMIT 20""",
            (session["user"],)
        ).fetchall()
    return jsonify([dict(r) for r in rows])


# ── ADMIN FRONTEND ─────────────────────────────────────────────────────────────

@app.route("/admin")
def admin_page():
    return send_from_directory("frontend", "admin.html")


# ── SERVE FRONTEND ─────────────────────────────────────────────────────────────

@app.route("/")
def index():
    return send_from_directory("frontend", "index.html")


# ── START ──────────────────────────────────────────────────────────────────────

if __name__ == "__main__":
    init_db()
    print("\n" + "="*55)
    print("  StartupScale360 — DD Intelligence Platform")
    print(f"  App:   http://localhost:5000")
    print(f"  Admin: http://localhost:5000/admin")
    print(f"  Admin password: {ADMIN_PASSWORD}")
    print("="*55 + "\n")
    port = int(os.environ.get("PORT", 5000))
    app.run(debug=False, host="0.0.0.0", port=port, threaded=True)